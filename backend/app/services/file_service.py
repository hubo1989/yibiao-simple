"""文件处理服务"""
import aiofiles
import os
import time
import gc
import hashlib
import io
import base64
import logging
from datetime import datetime
from typing import Optional, List, Dict, Tuple, Any
import PyPDF2
import docx
from fastapi import UploadFile
from ..config import settings

logger = logging.getLogger(__name__)

# 新增的第三方库 - 延迟导入以避免启动时卡住
HAS_ADVANCED_LIBS = False
HAS_OCR = False
HAS_EXCEL = False
HAS_PPTX = False
pdfplumber = None
fitz = None
docx2python = None
Image = None
pytesseract = None
openpyxl = None
pptx_module = None

def _init_advanced_libs():
    """延迟初始化高级文档处理库"""
    global HAS_ADVANCED_LIBS, pdfplumber, fitz, docx2python, Image
    global HAS_OCR, pytesseract
    global HAS_EXCEL, openpyxl
    global HAS_PPTX, pptx_module
    if HAS_ADVANCED_LIBS:
        return True
    try:
        import importlib
        pdfplumber = importlib.import_module('pdfplumber')
        fitz = importlib.import_module('fitz')
        docx2python = importlib.import_module('docx2python').docx2python
        Image = importlib.import_module('PIL.Image')
        HAS_ADVANCED_LIBS = True
    except ImportError as e:
        HAS_ADVANCED_LIBS = False
        print(f"高级文档处理库未安装: {e}")

    # OCR 支持
    try:
        import importlib
        pytesseract = importlib.import_module('pytesseract')
        HAS_OCR = True
    except ImportError:
        HAS_OCR = False

    # Excel 支持
    try:
        import importlib
        openpyxl = importlib.import_module('openpyxl')
        HAS_EXCEL = True
    except ImportError:
        HAS_EXCEL = False

    # PPT 支持
    try:
        import importlib
        pptx_module = importlib.import_module('pptx')
        HAS_PPTX = True
    except ImportError:
        HAS_PPTX = False

    return HAS_ADVANCED_LIBS

# 模块加载时初始化可选依赖
_init_advanced_libs()


class FileService:
    """文件处理服务"""

    # 图片存储目录
    EXTRACTED_IMAGES_DIR = os.path.join(settings.upload_dir, "extracted_images")

    @staticmethod
    def save_extracted_image(
        image_data: bytes,
        ext: str = "jpg",
        doc_hash: str = "unknown",
        page: int = 0,
        index: int = 0,
    ) -> Tuple[str, int, int]:
        """保存提取的图片到磁盘，返回 (文件路径, 宽度, 高度)"""
        os.makedirs(FileService.EXTRACTED_IMAGES_DIR, exist_ok=True)
        filename = f"{doc_hash}_{page}_{index}.{ext}"
        file_path = os.path.join(FileService.EXTRACTED_IMAGES_DIR, filename)
        with open(file_path, "wb") as f:
            f.write(image_data)
        # 获取图片尺寸
        width, height = 0, 0
        try:
            if Image:
                with Image.open(io.BytesIO(image_data)) as img:
                    width, height = img.size
        except Exception:
            pass
        return file_path, width, height

    @staticmethod
    def _compute_doc_hash(file_path: str) -> str:
        """计算文件的短哈希，用于图片命名"""
        h = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()[:12]

    @staticmethod
    def image_to_data_uri(image_data: bytes, ext: str = "jpg") -> str:
        """[已弃用] 兼容旧调用，内部转为 save_extracted_image"""
        file_path, width, height = FileService.save_extracted_image(image_data, ext)
        return file_path

    @staticmethod
    def extract_images_from_pdf(file_path: str) -> List[Tuple[bytes, str, int, int]]:
        """从PDF提取图片，返回 (图片数据, 扩展名, 页码, 图片索引) 列表"""
        if not HAS_ADVANCED_LIBS:
            return []

        images = []
        try:
            doc = fitz.open(file_path)

            for page_num in range(doc.page_count):
                page = doc[page_num]
                image_list = page.get_images(full=True)

                for img_index, img in enumerate(image_list):
                    try:
                        # 获取图片数据
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)

                        # 转换为RGB格式（如果是CMYK）
                        if pix.n - pix.alpha < 4:
                            img_data = pix.tobytes("jpeg")
                            ext = "jpg"
                        else:
                            pix1 = fitz.Pixmap(fitz.csRGB, pix)
                            img_data = pix1.tobytes("jpeg")
                            ext = "jpg"
                            pix1 = None

                        pix = None
                        images.append((img_data, ext, page_num + 1, img_index + 1))

                    except Exception as e:
                        print(f"提取PDF第{page_num+1}页图片{img_index+1}失败: {str(e)}")
                        continue

            doc.close()
            return images

        except Exception as e:
            print(f"PDF图片提取失败: {str(e)}")
            return []

    @staticmethod
    def extract_images_from_docx(file_path: str) -> List[Tuple[bytes, str, int]]:
        """从Word文档提取图片，返回 (图片数据, 扩展名, 图片索引) 列表"""
        images = []
        doc = None
        try:
            doc = docx.Document(file_path)

            # 获取文档中的所有关系
            rels = doc.part.rels
            img_index = 0

            for rel in rels.values():
                if "image" in rel.target_ref:
                    try:
                        # 读取图片数据
                        img_data = rel.target_part.blob

                        # 根据content_type确定扩展名
                        content_type = rel.target_part.content_type
                        if 'jpeg' in content_type:
                            ext = 'jpg'
                        elif 'png' in content_type:
                            ext = 'png'
                        elif 'gif' in content_type:
                            ext = 'gif'
                        elif 'bmp' in content_type:
                            ext = 'bmp'
                        else:
                            ext = 'jpg'  # 默认

                        img_index += 1
                        images.append((img_data, ext, img_index))

                    except Exception as e:
                        print(f"提取Word文档图片{img_index+1}失败: {str(e)}")
                        continue

            if doc:
                del doc
            gc.collect()
            return images

        except Exception as e:
            if doc:
                del doc
            gc.collect()
            print(f"Word文档图片提取失败: {str(e)}")
            return []

    @staticmethod
    def _safe_file_cleanup(file_path: str, max_retries: int = 3) -> bool:
        """安全删除文件，带重试机制"""
        for attempt in range(max_retries):
            try:
                if os.path.exists(file_path):
                    # 强制垃圾回收，释放可能的文件句柄
                    gc.collect()
                    time.sleep(0.1 * (attempt + 1))  # 递增延迟
                    os.remove(file_path)
                return True
            except OSError as e:
                if attempt == max_retries - 1:
                    print(f"无法删除文件 {file_path}: {e}")
                    return False
                time.sleep(0.5)  # 等待后重试
        return True
    
    @staticmethod
    async def save_uploaded_file(file: UploadFile) -> str:
        """保存上传的文件并返回文件路径"""
        # 创建上传目录
        os.makedirs(settings.upload_dir, exist_ok=True)

        # 生成带时间戳的文件名，防止重复
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:-3]  # 精确到毫秒
        filename = file.filename or "unknown_file"

        # 分离文件名和扩展名
        name, ext = os.path.splitext(filename)

        # 生成新的文件名：原文件名_时间戳.扩展名
        new_filename = f"{name}_{timestamp}{ext}"
        file_path = os.path.join(settings.upload_dir, new_filename)

        # 异步保存文件
        async with aiofiles.open(file_path, 'wb') as f:
            content = await file.read()
            await f.write(content)

        return file_path
    
    @staticmethod
    async def extract_text_from_pdf(file_path: str) -> str:
        """从PDF文件提取文本，支持表格内容和图片。

        降级策略:
        1. pdfplumber (含表格 + 图片)
        2. PyMuPDF (pdfplumber 异常时自动回退)
        3. PyPDF2 (最基础)
        4. OCR (前三层文本 < 100 字符时触发，疑似扫描件)
        """
        import asyncio
        text = ""
        if HAS_ADVANCED_LIBS:
            text = await FileService._extract_pdf_with_pdfplumber(file_path)
        else:
            text = await asyncio.get_event_loop().run_in_executor(
                None, FileService._extract_pdf_with_pypdf2, file_path
            )

        # OCR 降级: 前三层提取文本过少，可能是扫描件
        if len(text) < 100 and HAS_OCR:
            logger.warning(
                "PDF 文本提取结果不足 100 字符 (%d)，尝试 OCR 降级: %s",
                len(text), file_path,
            )
            try:
                ocr_text = await asyncio.get_event_loop().run_in_executor(
                    None, FileService._extract_pdf_with_ocr, file_path
                )
                if ocr_text and len(ocr_text) > len(text):
                    text = ocr_text
            except Exception as e:
                logger.warning("OCR 降级失败，保留已有文本: %s", str(e))

        return text
    
    @staticmethod
    async def _extract_pdf_with_pdfplumber(file_path: str) -> str:
        """使用pdfplumber提取PDF文本，包含表格和图片（确保及时释放文件句柄）"""
        try:
            extracted_text = []
            image_references = []  # 存储图片文件路径映射
            global_img_counter = 1
            doc_hash = FileService._compute_doc_hash(file_path)

            # 获取PDF文档的所有图片信息，用于后续匹配
            all_images = FileService.extract_images_from_pdf(file_path)
            page_images_map = {}
            for img_data, ext, page_num, img_index in all_images:
                if page_num not in page_images_map:
                    page_images_map[page_num] = []
                page_images_map[page_num].append((img_data, ext, img_index))

            # 使用上下文管理器，避免在Windows上产生文件锁
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # 添加页码标识
                    extracted_text.append(f"\n--- 第 {page_num} 页 ---\n")

                    # 提取普通文本
                    text = page.extract_text()
                    if text:
                        # 检查文本中是否有图片标记
                        import re
                        img_pattern = r'----.*?(?:image|img|media).*?----'
                        img_matches = list(re.finditer(img_pattern, text, re.IGNORECASE))

                        if img_matches and page_num in page_images_map:
                            # 按顺序处理页面中的图片
                            page_images = page_images_map[page_num]
                            processed_text = text

                            for i, match in enumerate(img_matches):
                                if i < len(page_images):
                                    img_data, ext, img_index = page_images[i]

                                    # 保存图片到磁盘，生成元数据占位符
                                    saved_path, w, h = FileService.save_extracted_image(
                                        img_data, ext, doc_hash, page_num, global_img_counter
                                    )

                                    old_mark = match.group()
                                    new_mark = f"[图片{global_img_counter}: 位于第{page_num}页, {w}x{h}px]"
                                    processed_text = processed_text.replace(old_mark, new_mark, 1)

                                    image_references.append(f"[图片{global_img_counter}]: {saved_path}")
                                    global_img_counter += 1

                            extracted_text.append(processed_text)
                        else:
                            extracted_text.append(text)

                    # 提取表格
                    tables = page.extract_tables()
                    for table_num, table in enumerate(tables, 1):
                        extracted_text.append(f"\n[表格 {table_num}]")
                        for row in table:
                            if row:  # 跳过空行
                                # 过滤空值并连接单元格
                                row_text = " | ".join([str(cell) if cell else "" for cell in row])
                                extracted_text.append(row_text)
                        extracted_text.append("[表格结束]\n")

            # 在文档末尾添加图片文件路径映射（供前端按需加载）
            if image_references:
                extracted_text.append(f"\n\n--- 图片引用 ---")
                extracted_text.extend(image_references)

            result = "\n".join(extracted_text).strip()
            gc.collect()
            return result
        except Exception as e:
            gc.collect()
            # 如果pdfplumber失败，尝试PyMuPDF
            try:
                return await FileService._extract_pdf_with_pymupdf(file_path)
            except Exception:
                raise Exception(f"PDF文件读取失败: {str(e)}")
    
    @staticmethod
    async def _extract_pdf_with_pymupdf(file_path: str) -> str:
        """使用PyMuPDF提取PDF文本和图片"""
        import asyncio
        return await asyncio.get_event_loop().run_in_executor(
            None, FileService._extract_pdf_with_pymupdf_sync, file_path
        )

    @staticmethod
    def _extract_pdf_with_pymupdf_sync(file_path: str) -> str:
        """同步使用PyMuPDF提取PDF文本和图片"""
        try:
            doc = fitz.open(file_path)
            extracted_text = []
            
            for page_num in range(doc.page_count):
                page = doc[page_num]
                extracted_text.append(f"\n--- 第 {page_num + 1} 页 ---\n")
                
                # 提取文本
                text = page.get_text()
                if text:
                    extracted_text.append(text)
                
                # 尝试提取表格
                try:
                    tables = page.find_tables()
                    for table_num, table in enumerate(tables, 1):
                        extracted_text.append(f"\n[表格 {table_num}]")
                        table_data = table.extract()
                        for row in table_data:
                            if row:
                                row_text = " | ".join([str(cell) if cell else "" for cell in row])
                                extracted_text.append(row_text)
                        extracted_text.append("[表格结束]\n")
                except:
                    # 如果表格提取失败，跳过
                    pass
            
            doc.close()
            return "\n".join(extracted_text).strip()
        except Exception as e:
            raise Exception(f"PDF文件读取失败: {str(e)}")
    
    @staticmethod 
    def _extract_pdf_with_pypdf2(file_path: str) -> str:
        """使用PyPDF2提取PDF文本（原方法）"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text.strip()
        except Exception as e:
            raise Exception(f"PDF文件读取失败: {str(e)}")

    @staticmethod
    def _extract_pdf_with_ocr(file_path: str) -> str:
        """OCR 降级: 用 PyMuPDF 渲染每页为图片，再用 tesseract OCR 提取文字"""
        if not HAS_OCR or not fitz:
            raise Exception("OCR 依赖不可用")
        try:
            doc = fitz.open(file_path)
            extracted_text = []
            for page_num in range(doc.page_count):
                page = doc[page_num]
                extracted_text.append(f"\n--- 第 {page_num + 1} 页 (OCR) ---\n")
                # 渲染为 300 dpi 图片
                mat = fitz.Matrix(300 / 72, 300 / 72)
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")
                pix = None
                # OCR
                pil_img = Image.open(io.BytesIO(img_data))
                ocr_text = pytesseract.image_to_string(pil_img, lang="chi_sim+eng")
                if ocr_text and ocr_text.strip():
                    extracted_text.append(ocr_text.strip())
            doc.close()
            return "\n".join(extracted_text).strip()
        except Exception as e:
            raise Exception(f"PDF OCR 提取失败: {str(e)}")
    
    @staticmethod
    async def extract_text_from_docx(file_path: str) -> str:
        """从Word文档提取文本，支持表格内容和图片"""
        if HAS_ADVANCED_LIBS:
            return await FileService._extract_docx_with_docx2python(file_path)
        else:
            # 降级到原来的python-docx方法，但增强表格处理
            return await FileService._extract_docx_with_python_docx(file_path)
    
    @staticmethod
    async def _extract_docx_with_docx2python(file_path: str) -> str:
        """使用docx2python提取Word文档内容和图片（确保及时释放文件句柄）"""
        try:
            extracted_text = []
            image_references = []  # 存储图片文件路径映射
            global_img_counter = 1
            doc_hash = FileService._compute_doc_hash(file_path)

            # 获取Word文档的所有图片信息
            all_images = FileService.extract_images_from_docx(file_path)

            # 使用上下文管理器确保文件及时关闭，避免Windows上的锁定
            with docx2python(file_path) as content:
                # 处理文档内容
                if hasattr(content, 'document'):
                    for section in content.document:
                        for element in section:
                            if isinstance(element, list):
                                # 这可能是表格
                                extracted_text.append("\n[表格内容]")
                                for row in element:
                                    if isinstance(row, list):
                                        row_text = " | ".join([str(cell).strip() for cell in row if cell])
                                        if row_text:
                                            extracted_text.append(row_text)
                                    else:
                                        extracted_text.append(str(row))
                                extracted_text.append("[表格结束]\n")
                            else:
                                # 普通文本，检查是否包含图片标记
                                text = str(element).strip()
                                if text:
                                    # 检查文本中是否有图片标记
                                    import re
                                    img_pattern = r'----.*?(?:image|img|media).*?----'
                                    img_matches = list(re.finditer(img_pattern, text, re.IGNORECASE))

                                    if img_matches and all_images:
                                        processed_text = text

                                        for match in img_matches:
                                            if global_img_counter <= len(all_images):
                                                img_data, ext, img_index = all_images[global_img_counter - 1]

                                                # 保存图片到磁盘
                                                saved_path, w, h = FileService.save_extracted_image(
                                                    img_data, ext, doc_hash, 0, global_img_counter
                                                )

                                                old_mark = match.group()
                                                new_mark = f"[图片{global_img_counter}: {w}x{h}px]"
                                                processed_text = processed_text.replace(old_mark, new_mark, 1)

                                                image_references.append(f"[图片{global_img_counter}]: {saved_path}")
                                                global_img_counter += 1

                                        extracted_text.append(processed_text)
                                    else:
                                        extracted_text.append(text)

            # 在文档末尾添加图片文件路径映射
            if image_references:
                extracted_text.append(f"\n\n--- 图片引用 ---")
                extracted_text.extend(image_references)

            result = "\n".join(extracted_text).strip()
            gc.collect()
            return result
        except Exception as e:
            gc.collect()
            # 如果docx2python失败，回退到增强的python-docx
            try:
                return await FileService._extract_docx_with_python_docx(file_path)
            except Exception:
                raise Exception(f"Word文档读取失败: {str(e)}")
    
    @staticmethod
    async def _extract_docx_with_python_docx(file_path: str) -> str:
        """使用python-docx提取Word文档内容和图片（增强版）"""
        doc = None
        try:
            doc = docx.Document(file_path)
            extracted_text = []
            image_references = []  # 存储图片文件路径映射
            global_img_counter = 1
            doc_hash = FileService._compute_doc_hash(file_path)

            # 获取Word文档的所有图片信息
            all_images = FileService.extract_images_from_docx(file_path)

            # 提取段落文本，同时处理图片
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # 检查文本中是否有图片标记
                    import re
                    img_pattern = r'----.*?(?:image|img|media).*?----'
                    img_matches = list(re.finditer(img_pattern, text, re.IGNORECASE))

                    if img_matches and all_images:
                        processed_text = text

                        for match in img_matches:
                            if global_img_counter <= len(all_images):
                                img_data, ext, img_index = all_images[global_img_counter - 1]

                                # 保存图片到磁盘
                                saved_path, w, h = FileService.save_extracted_image(
                                    img_data, ext, doc_hash, 0, global_img_counter
                                )

                                old_mark = match.group()
                                new_mark = f"[图片{global_img_counter}: {w}x{h}px]"
                                processed_text = processed_text.replace(old_mark, new_mark, 1)

                                image_references.append(f"[图片{global_img_counter}]: {saved_path}")
                                global_img_counter += 1

                        extracted_text.append(processed_text)
                    else:
                        extracted_text.append(text)

            # 提取表格内容
            for table_num, table in enumerate(doc.tables, 1):
                extracted_text.append(f"\n[表格 {table_num}]")
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        row_data.append(cell_text if cell_text else "")
                    row_text = " | ".join(row_data)
                    if row_text.strip():
                        extracted_text.append(row_text)
                extracted_text.append("[表格结束]\n")

            # 在文档末尾添加图片文件路径映射
            if image_references:
                extracted_text.append(f"\n\n--- 图片引用 ---")
                extracted_text.extend(image_references)

            result = "\n".join(extracted_text).strip()

            # 确保释放资源
            if doc:
                del doc
            gc.collect()

            return result
        except Exception as e:
            # 确保释放资源
            if doc:
                del doc
            gc.collect()
            raise Exception(f"Word文档读取失败: {str(e)}")
    
    @staticmethod
    async def extract_text_from_excel(file_path: str) -> str:
        """从 Excel (.xlsx) 文件提取文本"""
        if not HAS_EXCEL:
            raise Exception("Excel 解析依赖 openpyxl 未安装")
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            extracted_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                extracted_text.append(f"\n--- Sheet: {sheet_name} ---\n")
                # 获取合并单元格信息
                merged_cells_map: Dict[Tuple[int, int], str] = {}
                if hasattr(ws, 'merged_cells') and ws.merged_cells:
                    for merged_range in ws.merged_cells.ranges:
                        top_left_value = None
                        for row_idx in range(merged_range.min_row, merged_range.max_row + 1):
                            for col_idx in range(merged_range.min_col, merged_range.max_col + 1):
                                if row_idx == merged_range.min_row and col_idx == merged_range.min_col:
                                    cell = ws.cell(row=row_idx, column=col_idx)
                                    top_left_value = str(cell.value) if cell.value is not None else "-"
                                else:
                                    merged_cells_map[(row_idx, col_idx)] = top_left_value or "-"

                for row in ws.iter_rows():
                    cells = []
                    for cell in row:
                        # 检查是否是合并单元格的非主单元格
                        key = (cell.row, cell.column)
                        if key in merged_cells_map:
                            cells.append(merged_cells_map[key])
                        elif cell.value is not None:
                            cells.append(str(cell.value))
                        else:
                            cells.append("-")
                    row_text = "\t".join(cells)
                    extracted_text.append(row_text)
            wb.close()
            return "\n".join(extracted_text).strip()
        except Exception as e:
            raise Exception(f"Excel 文件读取失败: {str(e)}")

    @staticmethod
    async def extract_text_from_pptx(file_path: str) -> str:
        """从 PowerPoint (.pptx) 文件提取文本"""
        if not HAS_PPTX:
            raise Exception("PPT 解析依赖 python-pptx 未安装")
        try:
            from pptx.util import Inches  # noqa: F401
            prs = pptx_module.Presentation(file_path)
            extracted_text = []
            img_counter = 1

            for slide_num, slide in enumerate(prs.slides, 1):
                # 提取幻灯片标题
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                extracted_text.append(f"\n--- Slide {slide_num}: {title or '(无标题)'} ---\n")

                for shape in slide.shapes:
                    # 文本框
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                extracted_text.append(para_text)

                    # 表格
                    if shape.has_table:
                        table = shape.table
                        extracted_text.append("\n[表格]")
                        for row in table.rows:
                            cells = []
                            for cell in row.cells:
                                cell_text = cell.text.strip() if cell.text else "-"
                                cells.append(cell_text)
                            extracted_text.append("\t".join(cells))
                        extracted_text.append("[表格结束]\n")

                    # 图片占位符
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        extracted_text.append(f"[图片{img_counter}: 幻灯片{slide_num}]")
                        img_counter += 1

                # 备注
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        extracted_text.append(f"\n[备注] {notes_text}")

            return "\n".join(extracted_text).strip()
        except Exception as e:
            raise Exception(f"PPT 文件读取失败: {str(e)}")

    @staticmethod
    async def process_uploaded_file(file: UploadFile) -> str:
        """处理上传的文件并提取文本内容"""
        # 检查文件大小
        content = await file.read()
        if len(content) > settings.max_file_size:
            raise Exception(f"文件大小超过限制 ({settings.max_file_size / 1024 / 1024}MB)")

        # 重置文件指针
        await file.seek(0)

        # 保存文件
        file_path = await FileService.save_uploaded_file(file)

        try:
            # 根据文件类型提取文本和图片
            content_type = file.content_type or ""
            if content_type == "application/pdf":
                text = await FileService.extract_text_from_pdf(file_path)
            elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                text = await FileService.extract_text_from_docx(file_path)
            elif content_type in (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
            ):
                text = await FileService.extract_text_from_excel(file_path)
            elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                text = await FileService.extract_text_from_pptx(file_path)
            else:
                raise Exception("不支持的文件类型，请上传PDF、Word、Excel或PPT文档")

            # 成功提取后，使用安全的文件清理方法
            FileService._safe_file_cleanup(file_path)

            return text

        except Exception as e:
            # 异常情况下也使用安全的文件清理方法
            FileService._safe_file_cleanup(file_path)
            raise e

    # ============ 素材解析：文本与图片分离 ============

    @staticmethod
    async def _save_material_image(image_data: bytes, ext: str, material_id: str, index: int) -> str:
        """将单张素材图片保存到磁盘（async），返回相对路径"""
        images_dir = os.path.join(settings.upload_dir, "material_images", material_id)
        os.makedirs(images_dir, exist_ok=True)
        filename = f"image_{index:03d}.{ext}"
        file_path = os.path.join(images_dir, filename)
        async with aiofiles.open(file_path, "wb") as f:
            await f.write(image_data)
        # 返回相对于 upload_dir 的路径，方便后续 API 暴露
        return os.path.join("material_images", material_id, filename)

    @staticmethod
    def _safe_image_dir_cleanup(material_id: str) -> None:
        """删除 material_images/{material_id}/ 目录及其所有内容"""
        import shutil
        images_dir = os.path.join(settings.upload_dir, "material_images", material_id)
        try:
            if os.path.isdir(images_dir):
                shutil.rmtree(images_dir, ignore_errors=True)
        except Exception as e:
            print(f"清理素材图片目录失败 {images_dir}: {e}")

    @staticmethod
    async def _extract_pdf_with_images(file_path: str, material_id: str) -> Dict[str, Any]:
        """从 PDF 提取文本和图片（分离模式），返回 {text, images}"""
        extracted_text_parts: List[str] = []
        images: List[Dict[str, Any]] = []
        global_img_counter = 1

        # 先提取所有图片并保存到磁盘
        all_raw_images = FileService.extract_images_from_pdf(file_path)
        page_images_map: Dict[int, List[int]] = {}  # page_num -> [img_index, ...]
        for img_data, ext, page_num, img_index in all_raw_images:
            rel_path = await FileService._save_material_image(img_data, ext, material_id, global_img_counter)
            images.append({
                "index": global_img_counter,
                "marker": f"[图片{global_img_counter}]",
                "file_path": rel_path,
                "page": page_num,
                "format": ext,
                "size": len(img_data),
            })
            if page_num not in page_images_map:
                page_images_map[page_num] = []
            page_images_map[page_num].append(global_img_counter)
            global_img_counter += 1

        # 再提取文本，并在每页末尾追加该页的图片标记
        try:
            if pdfplumber is not None:
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        extracted_text_parts.append(f"\n--- 第 {page_num} 页 ---\n")
                        text = page.extract_text()
                        if text:
                            extracted_text_parts.append(text)
                        tables = page.extract_tables()
                        for table_num, table in enumerate(tables, 1):
                            extracted_text_parts.append(f"\n[表格 {table_num}]")
                            for row in table:
                                if row:
                                    row_text = " | ".join([str(cell) if cell else "" for cell in row])
                                    extracted_text_parts.append(row_text)
                            extracted_text_parts.append("[表格结束]\n")
                        # 在页文本末尾追加本页图片标记
                        for idx in page_images_map.get(page_num, []):
                            extracted_text_parts.append(f"[图片{idx}]")
            elif fitz is not None:
                doc = fitz.open(file_path)
                for page_num in range(doc.page_count):
                    page = doc[page_num]
                    extracted_text_parts.append(f"\n--- 第 {page_num + 1} 页 ---\n")
                    text = page.get_text()
                    if text:
                        extracted_text_parts.append(text)
                    # 在页文本末尾追加本页图片标记
                    for idx in page_images_map.get(page_num + 1, []):
                        extracted_text_parts.append(f"[图片{idx}]")
                doc.close()
            else:
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page_no, page in enumerate(pdf_reader.pages, 1):
                        extracted_text_parts.append(page.extract_text() or "")
                        for idx in page_images_map.get(page_no, []):
                            extracted_text_parts.append(f"[图片{idx}]")
        except Exception as e:
            raise Exception(f"PDF文本提取失败: {str(e)}")

        text = "\n".join(extracted_text_parts).strip()
        return {"text": text, "images": images}

    @staticmethod
    async def _extract_docx_with_images(file_path: str, material_id: str) -> Dict[str, Any]:
        """从 DOCX 提取文本和图片（分离模式），返回 {text, images}"""
        images: List[Dict[str, Any]] = []
        global_img_counter = 1

        # 先提取所有图片并保存到磁盘，建立 relationship id → 图片序号 的映射
        all_raw_images = FileService.extract_images_from_docx(file_path)
        for img_data, ext, img_index in all_raw_images:
            rel_path = await FileService._save_material_image(img_data, ext, material_id, global_img_counter)
            images.append({
                "index": global_img_counter,
                "marker": f"[图片{global_img_counter}]",
                "file_path": rel_path,
                "format": ext,
                "size": len(img_data),
            })
            global_img_counter += 1

        # 再提取文本，通过检测段落中的图片 run 注入 [图片N] 标记
        extracted_text_parts: List[str] = []
        inline_img_counter = 1
        try:
            import lxml.etree as etree  # python-docx 依赖 lxml，已安装
            doc = docx.Document(file_path)

            # 收集文档中图片所在的段落位置（用 XML 检测 <a:blip> 或 <v:imagedata>）
            _NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
            _NS_V = "urn:schemas-microsoft-com:vml"
            _NS_W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

            def _paragraph_has_image(para) -> bool:
                """检测段落 XML 中是否包含内联图片标记"""
                xml_str = para._element.xml
                return (
                    "a:blip" in xml_str
                    or "v:imagedata" in xml_str
                    or "pic:pic" in xml_str
                    or ('<w:drawing' in xml_str)
                    or ('<w:pict' in xml_str)
                )

            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if _paragraph_has_image(paragraph) and inline_img_counter <= len(images):
                    # 先输出段落原有文字（如有），再插入图片标记
                    if text:
                        extracted_text_parts.append(text)
                    extracted_text_parts.append(f"[图片{inline_img_counter}]")
                    inline_img_counter += 1
                elif text:
                    extracted_text_parts.append(text)

            for table_num, table in enumerate(doc.tables, 1):
                extracted_text_parts.append(f"\n[表格 {table_num}]")
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        row_data.append(cell_text if cell_text else "")
                    row_text = " | ".join(row_data)
                    if row_text.strip():
                        extracted_text_parts.append(row_text)
                extracted_text_parts.append("[表格结束]\n")

            # 如果有图片但段落检测未注入全部（fallback：在文末追加剩余标记）
            while inline_img_counter <= len(images):
                extracted_text_parts.append(f"[图片{inline_img_counter}]")
                inline_img_counter += 1

            del doc
            gc.collect()
        except Exception as e:
            gc.collect()
            raise Exception(f"Word文档文本提取失败: {str(e)}")

        text = "\n".join(extracted_text_parts).strip()
        return {"text": text, "images": images}

    @staticmethod
    async def parse_material(file: UploadFile) -> Dict[str, Any]:
        """
        解析上传的素材文件，将文本和图片分离返回。

        返回:
            {
                "text": str,           # 提取的纯文本内容
                "images": [            # 提取的图片列表
                    {
                        "index": int,
                        "marker": str,      # 如 "[图片1]"
                        "file_path": str,   # 图片文件相对路径
                        "format": str,      # 图片格式 (jpg/png/...)
                        "size": int,        # 图片字节大小
                        "page": int|None,   # PDF 页码（仅 PDF）
                    }
                ],
                "material_id": str,     # 本次解析的唯一 ID
                "source_filename": str, # 原始文件名
            }
        """
        content = await file.read()
        if len(content) > settings.max_file_size:
            raise Exception(f"文件大小超过限制 ({settings.max_file_size / 1024 / 1024}MB)")

        await file.seek(0)

        # 检查文件类型
        is_pdf = content[:5] == b'%PDF-'
        is_docx = content[:4] == b'PK\x03\x04'
        if not (is_pdf or is_docx):
            raise Exception("不支持的文件类型，请上传PDF或Word文档")

        material_id = uuid_mod.uuid4().hex[:12]
        file_path = await FileService.save_uploaded_file(file)
        extraction_ok = False

        try:
            _init_advanced_libs()
            if is_pdf:
                result = await FileService._extract_pdf_with_images(file_path, material_id)
            else:
                result = await FileService._extract_docx_with_images(file_path, material_id)

            result["material_id"] = material_id
            result["source_filename"] = file.filename or "unknown"
            extraction_ok = True
            return result
        finally:
            # 无论成功失败，都清理临时上传文件
            FileService._safe_file_cleanup(file_path)
            # 提取失败时，清理已写出的素材图片目录
            if not extraction_ok:
                FileService._safe_image_dir_cleanup(material_id)

    @staticmethod
    async def write_material_owner(material_id: str, owner_id: str) -> None:
        """将所有者 ID 写入 material_images/{material_id}/.owner 文件"""
        images_dir = os.path.join(settings.upload_dir, "material_images", material_id)
        os.makedirs(images_dir, exist_ok=True)
        owner_file = os.path.join(images_dir, ".owner")
        async with aiofiles.open(owner_file, "w") as f:
            await f.write(str(owner_id))

    @staticmethod
    def read_material_owner(material_id: str) -> str | None:
        """读取素材图片目录的所有者 ID，不存在则返回 None"""
        owner_file = os.path.join(settings.upload_dir, "material_images", material_id, ".owner")
        try:
            with open(owner_file, "r") as f:
                return f.read().strip()
        except OSError:
            return None